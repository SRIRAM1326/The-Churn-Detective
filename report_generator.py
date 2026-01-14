import pandas as pd
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
import os

# 1. Load Data
df = pd.read_csv('retention_insights.csv')

def generate_pdf_report():
    print("Generating PDF Executive Summary...")
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", 'B', 16)
    
    # Header
    pdf.cell(190, 10, "The Churn Detective: Executive Retention Report", ln=True, align='C')
    pdf.ln(10)
    
    # Executive Summary Metrics
    pdf.set_font("Arial", 'B', 12)
    pdf.cell(190, 10, "Executive Summary Metrics", ln=True)
    pdf.set_font("Arial", '', 11)
    
    total_revenue_at_risk = df['Revenue_at_Risk'].sum()
    high_risk_count = len(df[df['Risk_Level'] == 'High'])
    avg_churn_prob = df['Churn_Probability'].mean()
    
    pdf.cell(190, 8, f"- Total Customers Analyzed: {len(df):,}", ln=True)
    pdf.cell(190, 8, f"- High Risk Customers Identified: {high_risk_count:,}", ln=True)
    pdf.cell(190, 8, f"- Total Monthly Revenue at Risk: ${total_revenue_at_risk:,.2f}", ln=True)
    pdf.cell(190, 8, f"- Average Churn Probability: {avg_churn_prob:.2f}%", ln=True)
    pdf.ln(10)
    
    # Top 5 High-Value Retention Targets
    pdf.set_font("Arial", 'B', 12)
    pdf.cell(190, 10, "Top 5 High-Value Retention Targets", ln=True)
    pdf.set_font("Arial", '', 10)
    
    top_targets = df[df['Risk_Level'] == 'High'].sort_values(by='Revenue_at_Risk', ascending=False).head(5)
    
    # Table Header
    pdf.set_fill_color(200, 220, 255)
    pdf.cell(40, 8, "Customer ID", 1, 0, 'C', True)
    pdf.cell(40, 8, "Rev at Risk", 1, 0, 'C', True)
    pdf.cell(110, 8, "Recommended Strategy", 1, 1, 'C', True)
    
    for _, row in top_targets.iterrows():
        pdf.cell(40, 8, f"{row['Customer_ID']}", 1)
        pdf.cell(40, 8, f"${row['Revenue_at_Risk']:,.2f}", 1)
        pdf.cell(110, 8, f"{row['Retention_Strategy'][:55]}...", 1, 1)
        
    pdf.output("Churn_Detective_Executive_Report.pdf")
    print("PDF Report Saved: 'Churn_Detective_Executive_Report.pdf'")

def generate_pptx_summary():
    print("Generating PPTX Slide...")
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
    title_tf = title_box.text_frame
    title_tf.text = "Churn Detective: High-Value Risk Analysis"
    
    # Metrics Box
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Key Metrics"
    
    p = tf.add_paragraph()
    p.text = f"Total Revenue at Risk: ${df['Revenue_at_Risk'].sum():,.2f}"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p2 = tf.add_paragraph()
    p2.text = f"High Risk High Value Customers: {len(df[(df['Risk_Level'] == 'High') & (df['Is_High_Value'] == 1)])}"
    p2.font.size = Pt(14)

    # Strategy Highlight
    strategy_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4.5), Inches(3))
    st_tf = strategy_box.text_frame
    st_tf.text = "Top Retention Strategy"
    p3 = st_tf.add_paragraph()
    p3.text = "Priority Outreach for High-Revenue Segments"
    p3.font.bold = True
    
    prs.save("Churn_Detective_Executive_Slide.pptx")
    print("PPTX Slide Saved: 'Churn_Detective_Executive_Slide.pptx'")

if __name__ == "__main__":
    generate_pdf_report()
    generate_pptx_summary()
